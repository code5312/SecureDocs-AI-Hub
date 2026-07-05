from pathlib import Path
from typing import BinaryIO

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from app.config.settings import Settings
from app.extraction.errors import ExtractionError, ExtractionErrorCode
from app.extraction.types import ExtractedSegment

_TEXT_MIME = {"text/plain", "text/markdown"}


def _check_chars(total: int, settings: Settings) -> None:
    if total > settings.extraction_max_characters:
        raise ExtractionError(ExtractionErrorCode.CHARACTER_LIMIT_EXCEEDED)


def extract_text(file_obj: BinaryIO, *, filename: str, mime_type: str, settings: Settings) -> list[ExtractedSegment]:
    suffix = Path(filename).suffix.lower().lstrip(".")
    if suffix in {"txt", "md"} or mime_type in _TEXT_MIME:
        return _extract_text_file(file_obj, settings)
    if suffix == "pdf" or mime_type == "application/pdf":
        return _extract_pdf(file_obj, settings)
    if suffix == "docx":
        return _extract_docx(file_obj, settings)
    if suffix == "pptx":
        return _extract_pptx(file_obj, settings)
    if suffix == "xlsx":
        return _extract_xlsx(file_obj, settings)
    raise ExtractionError(ExtractionErrorCode.UNSUPPORTED_FORMAT)


def _extract_text_file(file_obj: BinaryIO, settings: Settings) -> list[ExtractedSegment]:
    data = file_obj.read()
    if b"\x00" in data:
        raise ExtractionError(ExtractionErrorCode.PARSER_ERROR)
    try:
        text = data.decode("utf-8-sig")
    except UnicodeDecodeError as exc:
        raise ExtractionError(ExtractionErrorCode.PARSER_ERROR) from exc
    _check_chars(len(text), settings)
    return [ExtractedSegment(text=text)]


def _extract_pdf(file_obj: BinaryIO, settings: Settings) -> list[ExtractedSegment]:
    try:
        reader = PdfReader(file_obj)
        if reader.is_encrypted:
            raise ExtractionError(ExtractionErrorCode.ENCRYPTED_DOCUMENT)
        if len(reader.pages) > settings.extraction_max_pages:
            raise ExtractionError(ExtractionErrorCode.PAGE_LIMIT_EXCEEDED)
        segments: list[ExtractedSegment] = []
        total = 0
        for index, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            total += len(text)
            _check_chars(total, settings)
            segments.append(ExtractedSegment(text=text, page_number=index))
        return segments
    except ExtractionError:
        raise
    except Exception as exc:
        raise ExtractionError(ExtractionErrorCode.PARSER_ERROR) from exc


def _extract_docx(file_obj: BinaryIO, settings: Settings) -> list[ExtractedSegment]:
    try:
        doc = DocxDocument(file_obj)
        segments: list[ExtractedSegment] = []
        total = 0
        current_heading: str | None = None
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            style = paragraph.style.name if paragraph.style else ""
            if style.startswith("Heading"):
                current_heading = text[:255]
            total += len(text)
            _check_chars(total, settings)
            segments.append(ExtractedSegment(text=text, section_title=current_heading))
        for table in doc.tables:
            rows = ["\t".join(cell.text.strip() for cell in row.cells if cell.text.strip()) for row in table.rows]
            text = "\n".join(row for row in rows if row)
            if text:
                total += len(text)
                _check_chars(total, settings)
                segments.append(ExtractedSegment(text=text, section_title=current_heading))
        return segments
    except Exception as exc:
        raise ExtractionError(ExtractionErrorCode.PARSER_ERROR) from exc


def _extract_pptx(file_obj: BinaryIO, settings: Settings) -> list[ExtractedSegment]:
    try:
        prs = Presentation(file_obj)
        if len(prs.slides) > settings.extraction_max_slides:
            raise ExtractionError(ExtractionErrorCode.SLIDE_LIMIT_EXCEEDED)
        segments: list[ExtractedSegment] = []
        total = 0
        for index, slide in enumerate(prs.slides, start=1):
            texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            text = "\n".join(texts)
            total += len(text)
            _check_chars(total, settings)
            segments.append(ExtractedSegment(text=text, slide_number=index))
        return segments
    except ExtractionError:
        raise
    except Exception as exc:
        raise ExtractionError(ExtractionErrorCode.PARSER_ERROR) from exc


def _extract_xlsx(file_obj: BinaryIO, settings: Settings) -> list[ExtractedSegment]:
    try:
        wb = load_workbook(file_obj, read_only=True, data_only=True, keep_links=False)
        if len(wb.worksheets) > settings.extraction_max_sheets:
            raise ExtractionError(ExtractionErrorCode.SHEET_LIMIT_EXCEEDED)
        segments: list[ExtractedSegment] = []
        total = 0
        for ws in wb.worksheets:
            rows: list[str] = []
            row_start: int | None = None
            row_end: int | None = None
            for row_index, row in enumerate(ws.iter_rows(values_only=True), start=1):
                if row_index > settings.extraction_max_rows_per_sheet:
                    raise ExtractionError(ExtractionErrorCode.SHEET_LIMIT_EXCEEDED)
                values = [str(value) for value in row if value not in (None, "")]
                if not values:
                    continue
                row_start = row_start or row_index
                row_end = row_index
                rows.append("\t".join(values))
            text = "\n".join(rows)
            if text:
                total += len(text)
                _check_chars(total, settings)
                segments.append(ExtractedSegment(text=text, sheet_name=ws.title[:255], row_start=row_start, row_end=row_end))
        wb.close()
        return segments
    except ExtractionError:
        raise
    except Exception as exc:
        raise ExtractionError(ExtractionErrorCode.PARSER_ERROR) from exc
